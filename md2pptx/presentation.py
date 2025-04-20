from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Create a presentation
prs = Presentation()

# Background image path
background = "/home/allanabiud/Desktop/School/JKUAT/CBD/4.1/COMPUTER SECURITY AND CRYPTOGRAPHY/ASSIGNMENT/BACKGROUND.jpg"

# Define content (Title -> Details)
slides_content = [
    (
        "COMPUTER FRAUD AND PREVENTION TECHNIQUES",
        "Overview of various computer fraud techniques and how to prevent them.",
    ),
    (
        "Cracking",
        "Definition: Bypassing security protections in software.\nExample: Removing licensing restrictions from paid software.\nHow it happens: Decompiling programs, manipulating code.\nPrevention: Strong encryption, DRM tools.",
    ),
    (
        "Data Diddling",
        "Definition: Manipulating data before/during entry.\nExample: A bank employee alters account balances.\nHow it happens: Tampering with records, modifying grades.\nPrevention: Input validation, encryption, audits.",
    ),
    (
        "Data Leakage",
        "Definition: Unauthorized transfer of sensitive data.\nExample: Employee leaks confidential business strategies.\nHow it happens: USB copying, email leaks.\nPrevention: DLP tools, monitoring transfers.",
    ),
    (
        "Denial of Service (DoS) Attack",
        "Definition: Overloading a system with excessive traffic.\nExample: Hacker crashes a website via massive requests.\nHow it happens: Flooding servers with data.\nPrevention: Firewalls, CDNs, rate-limiting.",
    ),
    (
        "Eavesdropping",
        "Definition: Intercepting private communications to gather sensitive information.\nExample: Hackers use Wi-Fi sniffers to capture login credentials on public networks.\nHow it happens: Using spyware, wiretapping, or packet sniffers.\nPrevention: Use encrypted communication channels and avoid unsecured Wi-Fi.",
    ),
    (
        "Email Forgery and Threats",
        "Definition: Sending fake emails to mislead, scam, or threaten recipients.\nExample: A scammer pretends to be a bank and tricks customers into revealing login details.\nHow it happens: Attackers forge email headers.\nPrevention: Verify email authenticity, use spam filters, educate users.",
    ),
    (
        "Hacking",
        "Definition: Gaining unauthorized access to a system.\nExample: A hacker breaks into a government database.\nHow it happens: Exploiting software vulnerabilities or stolen passwords.\nPrevention: Implement MFA, update software, monitor activity.",
    ),
    (
        "Internet Misinformation & Cyber Terrorism",
        "Definition: Spreading false information or using online platforms for crimes.\nExample: Fake news influences public opinion.\nHow it happens: Misleading articles, hacking, propaganda.\nPrevention: Verify sources, promote digital literacy, monitor threats.",
    ),
    (
        "Logic Time Bomb",
        "Definition: Malicious code that remains dormant until triggered.\nExample: A script deletes company data if a programmer is fired.\nHow it happens: Hidden malicious code activates at a set time.\nPrevention: Conduct security audits, scan for unauthorized scripts.",
    ),
    (
        "Masquerading (Impersonation)",
        "Definition: Pretending to be someone else to gain access.\nExample: A hacker logs in as an administrator using stolen credentials.\nHow it happens: Phishing, stolen passwords, social engineering.\nPrevention: Implement MFA, train employees.",
    ),
    (
        "Password Cracking",
        "Definition: Breaking weak passwords to gain access.\nExample: A brute-force attack cracks a short password.\nHow it happens: Automated tools guess passwords.\nPrevention: Use strong passwords, enable MFA.",
    ),
    (
        "Piggybacking",
        "Definition: Gaining unauthorized access by following an authorized user.\nExample: A stranger sneaks into a secured building.\nHow it happens: Tailgating, exploiting active login sessions.\nPrevention: Enforce badge access, train employees.",
    ),
    (
        "Round-Down & Salami Technique",
        "Definition: Stealing small amounts of money in unnoticed transactions.\nExample: Rounding down financial transactions and keeping the fractions.\nHow it happens: Manipulating calculations in banking/payroll systems.\nPrevention: Regular audits, strict financial controls.",
    ),
    (
        "Software Piracy",
        "Definition: Copying/distributing software without authorization.\nExample: Using cracked software instead of purchasing a license.\nPrevention: Implement anti-piracy measures, educate users.",
    ),
    (
        "Scavenging",
        "Definition: Retrieving discarded sensitive information.\nExample: A hacker finds login credentials in the trash.\nHow it happens: Dumpster diving, restoring deleted data.\nPrevention: Shred documents, use secure deletion tools.",
    ),
    (
        "Social Engineering",
        "Definition: Manipulating people into revealing confidential info.\nExample: A hacker pretends to be IT support.\nHow it happens: Phishing, pretexting, psychological manipulation.\nPrevention: Train employees, verify all requests.",
    ),
    (
        "Superzapping",
        "Definition: Using system privileges to bypass security controls.\nExample: An admin secretly accesses payroll data.\nHow it happens: Exploiting system privileges.\nPrevention: Restrict admin access, monitor user activity.",
    ),
    (
        "Trap Door (Backdoor)",
        "Definition: A hidden way to bypass authentication.\nExample: A developer leaves a secret admin login in software.\nHow it happens: Intentionally inserted in code but later exploited.\nPrevention: Conduct code reviews, remove unused access.",
    ),
    (
        "Trojan Horse, Virus, Worm",
        "Trojan Horse: Disguised as useful software but contains malware.\nVirus: Self-replicating malware that spreads through files.\nWorm: Malware that spreads across networks without user interaction.\nExample: A fake PDF file installs malware.\nPrevention: Use antivirus software, update security patches.",
    ),
]

# Deterring Computer Fraud (Slide 1)
deterring_fraud_content = (
    "• Strong Authentication & Access Controls – Use MFA, strong passwords, and role-based access.\n"
    "• Employee Training & Awareness – Educate on phishing, social engineering, and security best practices.\n"
    "• Security Policies & Procedures – Implement clear IT policies and regularly update them.\n"
    "• Data Encryption & Secure Communication – Encrypt sensitive data and use secure protocols (HTTPS, VPN).\n"
    "• Physical Security – Restrict access to systems using ID badges, biometrics, or security personnel.\n"
    "• Legal & Disciplinary Measures – Establish penalties and legal actions for fraud."
)

# Detecting Computer Fraud (Slide 2)
detecting_fraud_content = (
    "• Intrusion Detection & Firewalls – Monitor traffic and block unauthorized access.\n"
    "• Log Monitoring & Analysis – Track user activities and analyze logs for anomalies.\n"
    "• Fraud Detection Software – Use AI and machine learning to detect suspicious transactions.\n"
    "• Regular Audits & Penetration Testing – Assess security controls and identify weaknesses.\n"
    "• User Activity Monitoring – Detect abnormal behavior and unauthorized access attempts.\n"
    "• Reporting Mechanisms – Provide anonymous ways to report suspicious activities.\n"
    "• Backup & Recovery Plans – Maintain secure backups to prevent data loss from cyberattacks."
)


# Function to add slides
def add_slide(title, content):
    slide_layout = prs.slide_layouts[5]  # Title Only Layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background image
    slide.shapes.add_picture(
        background, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )

    # Add text box for title + content
    left = Inches(0.8)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Add title inside content box (Bold, Large, White)
    title_paragraph = text_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.bold = True
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    title_paragraph.font.name = "Archivo"
    title_paragraph.space_after = Pt(20)

    # Add content with bold labels
    for line in content.split("\n"):
        if ": " in line:  # Check if line has a label (e.g., "Definition:")
            label, text = line.split(": ", 1)
            p = text_frame.add_paragraph()
            run_label = p.add_run()
            run_label.text = label + ": "
            run_label.font.bold = True
            run_label.font.size = Pt(28)
            run_label.font.color.rgb = RGBColor(255, 255, 255)
            run_label.font.name = "Archivo"

            run_text = p.add_run()
            run_text.text = text
            run_text.font.size = Pt(28)
            run_text.font.color.rgb = RGBColor(255, 255, 255)
            run_text.font.name = "Archivo"
        else:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Archivo"


# Generate slides
for title, content in slides_content:
    add_slide(title, content)

# Add two slides
add_slide("Deterring Computer Fraud", deterring_fraud_content)
add_slide("Detecting Computer Fraud", detecting_fraud_content)

# Save the presentation
prs.save("computer_fraud.pptx")
print("Presentation saved as computer_fraud.pptx")
